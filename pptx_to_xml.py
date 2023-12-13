# -*- coding: utf-8 -*-
"""
Created on Mon Dec 11 14:11:18 2023

@author: rbuck
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_CATEGORY_TYPE
import xml.etree.ElementTree as ET
from datetime import datetime, timedelta

def serial_date_to_string(serial_date, date_format = '%m/%d/%Y'):
    """
    Converts serial date from Excel to real date as string

    Parameters
    ----------
    serial_date : int
        serial date from Excel.
    date_format : str, optional
        format for date. The default is '%m/%d/%Y'.

    Returns
    -------
    str
        date in month/day/year (or other specified format).

    """
    base_date = datetime(1899, 12, 30)
    offset = timedelta(days=serial_date)
    dt = base_date + offset
    return dt.strftime(date_format)

def get_element_type(element_type, enumerator):
    """
    Converts element number code to text name based on enumerator object

    Parameters
    ----------
    element_type : int
        number returned as type code from pptx element.
    enumerator : pptx.enum.base.MetaEnumeration
        MetaEnumeration object of element types from python-pptx.

    Returns
    -------
    str
        text name of element type.

    """
    for enum in enumerator.__members__:
        if enum.value == element_type:
            return enum.name
    

def ungroup_shapes(slide):
    """
    Iterates through all shapes on a pptx slide and adds the individual elements
    of each grouped object (deleting the group) while slide still contains 
    grouped objects

    Parameters
    ----------
    slide : pptx.slide.Slide
        python-pptx slide object.

    Returns
    -------
    None.

    """
    while any(shape.shape_type == MSO_SHAPE_TYPE.GROUP for shape in slide.shapes):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                group = shape.element
                parent = group.getparent()
                index = parent.index(group)
                for member in group:
                    parent.insert(index,member)
                    index+=1
                parent.remove(group)

def process_text(shape, slide_element):
    """
    Function to process pptx shape's text and add to xml tree 

    Parameters
    ----------
    shape : pptx.shape.Shape
        python-pptx shape object with text_frame.
    slide_element : xml.etree.ElementTree.Element
        parent slide in xml tree.

    Returns
    -------
    None.

    """
    if "Title" in shape.name:
        text_type = 'title'
    elif "Subtitle" in shape.name:
        text_type = 'subtitle'
    elif "Footnote" in shape.name:
        text_type = 'footnote'
    elif "Placeholder" in shape.name:
        text_type = 'placeholder'
    else:
        text_type = 'body'
    for paragraph in shape.text_frame.paragraphs:
        level = paragraph.level
        text_element = ET.SubElement(slide_element, "text", attrib={'text_type':text_type,'level':str(level),'shape_id':str(shape.shape_id),'shape_name':shape.name})
        text_element.text = paragraph.text
        
def process_table(shape, slide_element):
    """
    Function to process table in pptx slide and add to xml tree

    Parameters
    ----------
    shape : pptx.shape.Shape
        python-pptx shape object with table object.
    slide_element : xml.etree.ElementTree.Element
        parent slide in xml tree.

    Returns
    -------
    None.

    """
    table = shape.table
    n_rows = len(table.rows)
    n_cols = len(table.columns)
    table_element = ET.SubElement(slide_element, "table", attrib={'rows':str(n_rows),'columns':str(n_cols), 'shape_id': str(shape.shape_id)})
    for r in range(0, n_rows):
        for c in range(0,n_cols):
            cell = table.cell(r,c)
            cell_element = ET.SubElement(table_element, "cell", attrib={'row_id':str(r),'column_id':str(c)})
            if hasattr(cell, "text_frame"):
                for paragraph in cell.text_frame.paragraphs:
                    level = paragraph.level
                    text_element = ET.SubElement(cell_element, "text", attrib={'text_type':'table_text','level':str(level)})
                    text_element.text = paragraph.text

def process_chart(shape, slide_element):
    """
    Function to process chart in pptx slide and add to xml tree

    Parameters
    ----------
    shape : pptx.shape.Shape
        python-pptx shape object with chart object.
    slide_element : xml.etree.ElementTree.Element
        parent slide in xml tree.

    Returns
    -------
    None.

    """
    chart = shape.chart
    attributes = {}
    attributes['chart_type'] = get_element_type(chart.chart_type, XL_CHART_TYPE)
    attributes['plots'] = str(len(chart.plots))
    chart_element = ET.SubElement(slide_element,"chart", attrib = attributes)
    try:
        chart_title = chart.chart_title.text_frame.text
        if len(chart_title)>0:
            text_element = ET.SubElement(chart_element, "text", attrib={'text_type':'chart_title','level':'0'})
            text_element.text = chart_title
        else:
            pass
    except(ValueError):
        pass
    try:
        value_axis = chart.value_axis.axis_title.text_frame.text
        if len(value_axis)>0:
            text_element = ET.SubElement(chart_element, "text", attrib={'text_type':'value_axis_title','level':'0'})
            text_element.text = value_axis
        else:
            pass
    except(ValueError):
        pass
    try:
        category_axis = chart.category_axis.axis_title.text_frame.text
        if len(category_axis)>0:
            text_element = ET.SubElement(chart_element, "text", attrib={'text_type':'category_axis_title','level':'0'})
            text_element.text = category_axis
        else:
            pass
    except(ValueError):
        pass
    for plot in chart.plots:
        plot_element = ET.SubElement(chart_element,"plot", attrib={'series':str(len(plot.series))})
        for series in plot.series:
            series_element = ET.SubElement(plot_element, "series")
            text_element = ET.SubElement(series_element, "text", attrib={'text_type':'series_name','level':'0'})
            text_element.text = series.name
            data_labels = [value[0] for value in plot.categories.flattened_labels]
            try:
                if chart.category_axis.category_type == XL_CATEGORY_TYPE.TIME_SCALE:
                    data_labels = [serial_date_to_string(int(date)) for date in data_labels]
            except:
                pass
            column_element = ET.SubElement(series_element, "column", attrib={'axis': 'category_axis'})
            for point in data_labels:
                value_element = ET.SubElement(column_element, "value")
                value_element.text = str(point)
            column_element = ET.SubElement(series_element, "column", attrib={'axis': 'value_axis'})
            for point in series.values:
                value_element = ET.SubElement(column_element, "value")
                value_element.text = str(point)

def process_picture(shape, slide_element):
    """
    Function to process image in pptx slide adding image metadata and blob to xml

    Parameters
    ----------
    shape : pptx.shape.Shape
        python-pptx shape object of shape type Picture.
    slide_element : xml.etree.ElementTree.Element
        parent slide in xml tree.

    Returns
    -------
    None.

    """
    image = shape.image
    image_element = ET.SubElement(slide_element, "image", attrib={'image_type':image.content_type, 'dpi': str(image.dpi), 'size': str(image.size), 'shape_id': str(shape.shape_id)})
    image_element.text = image.filename
    blob_element = ET.SubElement(image_element, "blob")
    blob_element.text = str(image.blob)
    

def pptx_to_xml(pptx_path, xml_path = None, encoding = 'utf-8', pptx_attributes = None):
    """
    Extracts structured text, table, chart, and image data from Microsoft pptx 
    and saves as an xml document or returns an xml ElementTree if no path to 
    save is provided

    Parameters
    ----------
    pptx_path : str
        filepath of pptx.
    xml_path : str, optional
        filepath to save xml. The default is None.
    encoding : str, optional
        text encoding for saving xml. The default is 'utf-8'.
    pptx_attributes : dict, optional
        attribute information (metadata) for presentation. The default is None.

    Returns
    -------
    xml (xml.etree.ElementTree.ElementTree) if no path to save is provided.

    """
    powerpoint = Presentation(pptx_path)
    
    root = ET.Element('presentation')
    
    for i,slide in enumerate(powerpoint.slides):
        slide_element = ET.SubElement(root, "slide", attrib={'number':str(i), 'slide_id': str(slide.slide_id)})
                
        ungroup_shapes(slide)
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                if len(shape.text)>0:
                    process_text(shape, slide_element)
            if shape.has_table:
                process_table(shape, slide_element)
            if shape.has_chart:
                process_chart(shape, slide_element)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                process_picture(shape, slide_element)
                    
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if len(notes) > 0:
                text_element = ET.SubElement(slide_element, "text", attrib={'text_type':'note'})
                text_element.text = notes
            
    tree = ET.ElementTree(root)
    
    if pptx_attributes is not None:
        for key,value in pptx_attributes.items():
            if type(value) is not str:
                value = str(value)
            root.set(key,value)
    
    if xml_path is None:
        return(tree)
    else:
        tree.write(xml_path, encoding=encoding, xml_declaration= True)
